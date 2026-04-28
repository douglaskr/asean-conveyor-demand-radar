from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _add_title(slide, idx: int, title: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10.5), Inches(0.7))
    tf = title_box.text_frame
    tf.text = f"{idx}. {title}"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(14, 51, 87)


def _add_page_badge(slide, idx: int) -> None:
    badge = slide.shapes.add_textbox(Inches(11.2), Inches(0.2), Inches(1.8), Inches(0.6))
    tf = badge.text_frame
    tf.text = f"P{idx}/7"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(0, 105, 120)


def build_ppt_kr(texts: dict, charts: dict[str, Path], out_path: Path) -> Path:
    """Build KR deck without relying on layout placeholders."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for idx, title in enumerate(texts["pages"], start=1):
        slide = prs.slides.add_slide(blank_layout)
        _add_title(slide, idx, title)
        _add_page_badge(slide, idx)

        summary = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.2), Inches(1.1))
        summary.text_frame.text = "자동 생성 해석 박스: 점수와 뉴스 흐름을 기반으로 영업 우선순위를 제시합니다."

        if idx == 2 and charts.get("country") and charts["country"].exists():
            slide.shapes.add_picture(str(charts["country"]), Inches(0.7), Inches(2.2), width=Inches(11.5))
        if idx == 3 and charts.get("industry") and charts["industry"].exists():
            slide.shapes.add_picture(str(charts["industry"]), Inches(0.7), Inches(2.2), width=Inches(11.5))

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path
